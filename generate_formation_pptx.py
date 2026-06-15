"""Génération du PowerPoint de formation Suivi PDP/ICP - AGRIWATT
Date présentation : 01/06/2026
Auteur : T. HOCEDEZ
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# Couleurs Agriwatt
VERT = RGBColor(0x00, 0x69, 0x3E)
JAUNE = RGBColor(0xFF, 0xB8, 0x1C)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR = RGBColor(0xEE, 0xEE, 0xEE)
GRIS_MOYEN = RGBColor(0xBB, 0xBB, 0xBB)
GRIS_FONCE = RGBColor(0x55, 0x55, 0x55)
NOIR = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]


def set_text(tf, text, size=18, bold=False, color=NOIR, align=PP_ALIGN.LEFT, font="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                 color=NOIR, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    set_text(tf, text, size=size, bold=bold, color=color, align=align, font=font)
    return tb


def add_bullets(slide, left, top, width, height, items, size=18, color=NOIR,
                bullet_color=VERT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # puce
        rb = p.add_run()
        rb.text = "▸  "
        rb.font.name = "Calibri"
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = bullet_color
        # texte
        rt = p.add_run()
        rt.text = item
        rt.font.name = "Calibri"
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
    return tb


def add_header(slide, title):
    """Bandeau vert en haut + titre blanc."""
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.9))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = VERT
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.6),
                 title, size=24, bold=True, color=BLANC, anchor=MSO_ANCHOR.MIDDLE)
    # liseré jaune
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), SW, Inches(0.05))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = JAUNE


def add_footer(slide, idx, total):
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(8), Inches(0.3),
                 "Suivi PDP/ICP — Agriwatt — 01/06/2026", size=10, color=GRIS_FONCE)
    add_text_box(slide, Inches(11.3), Inches(7.1), Inches(1.5), Inches(0.3),
                 f"{idx} / {total}", size=10, color=GRIS_FONCE, align=PP_ALIGN.RIGHT)


def add_screenshot_placeholder(slide, left, top, width, height, description):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.line.color.rgb = GRIS_MOYEN
    sh.line.width = Pt(1.5)
    sh.fill.solid()
    sh.fill.fore_color.rgb = GRIS_CLAIR
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    set_text(tf, f"[ Capture à insérer ]\n{description}",
             size=14, bold=True, color=GRIS_FONCE, align=PP_ALIGN.CENTER)


def add_box(slide, left, top, width, height, title, body, fill=BLANC,
            border=VERT, title_color=VERT, body_color=NOIR, title_size=16, body_size=13):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.line.color.rgb = border
    sh.line.width = Pt(1.5)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Calibri"
    r1.font.size = Pt(title_size)
    r1.font.bold = True
    r1.font.color.rgb = title_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.line_spacing = 1.1
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Calibri"
    r2.font.size = Pt(body_size)
    r2.font.color.rgb = body_color


SCREENSHOTS = []  # liste de (slide_num, description)


def register_screenshot(slide_num, desc):
    SCREENSHOTS.append((slide_num, desc))


# ========================================================================
# Slide 1 — COUVERTURE
# ========================================================================
s = prs.slides.add_slide(BLANK)
# fond vert
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background()
bg.fill.solid()
bg.fill.fore_color.rgb = VERT
# bande jaune
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.2), SW, Inches(0.08))
b.line.fill.background()
b.fill.solid()
b.fill.fore_color.rgb = JAUNE
# logo placeholder
logo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5),
                          Inches(2.4), Inches(0.9))
logo.line.color.rgb = JAUNE
logo.line.width = Pt(1.5)
logo.fill.solid()
logo.fill.fore_color.rgb = VERT
tf = logo.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "AGRIWATT", size=20, bold=True, color=JAUNE, align=PP_ALIGN.CENTER)

add_text_box(s, Inches(0.6), Inches(2.3), Inches(12), Inches(1.2),
             "Suivi PDP / ICP", size=54, bold=True, color=BLANC)
add_text_box(s, Inches(0.6), Inches(3.4), Inches(12), Inches(0.8),
             "Application de pilotage HSE", size=28, color=JAUNE)
add_text_box(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.7),
             "Formation conducteurs de travaux", size=22, color=BLANC)
add_text_box(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.6),
             "01 juin 2026", size=20, bold=True, color=BLANC)
add_text_box(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             "T. HOCEDEZ — Responsable QSE", size=18, color=BLANC)

# ========================================================================
# Slide 2 — POURQUOI CETTE APP ?
# ========================================================================
TOTAL = 20  # on remplira après
s = prs.slides.add_slide(BLANK)
add_header(s, "Pourquoi cette application ?")
# 4 boîtes 2x2
box_w = Inches(5.8)
box_h = Inches(2.5)
gap = Inches(0.3)
left1 = Inches(0.6)
left2 = Inches(0.6) + box_w + gap
top1 = Inches(1.3)
top2 = top1 + box_h + gap
add_box(s, left1, top1, box_w, box_h,
        "Gain de temps",
        "Recherche Sellsy par n° OPP : adresse, GPS, contact client pré-remplis en 1 clic. "
        "Le ICP alimente automatiquement le PDP : zéro double-saisie.",
        title_size=20, body_size=15)
add_box(s, left2, top1, box_w, box_h,
        "Conformité légale",
        "Respect du Code du travail (R4511-1 et suivants). PDP, ICP et permis de feu "
        "signés et archivés : preuve juridique opposable en cas d'accident.",
        title_size=20, body_size=15)
add_box(s, left1, top2, box_w, box_h,
        "Mobilité terrain",
        "PWA installable sur smartphone / tablette. Mode 100 % hors-ligne, photos "
        "natives, GPS automatique, signature tactile.",
        title_size=20, body_size=15)
add_box(s, left2, top2, box_w, box_h,
        "Traçabilité",
        "Sync SharePoint multi-utilisateurs, historique complet, archivage automatique "
        "à la clôture, export PDF / Word / ZIP avec annexes.",
        title_size=20, body_size=15)
add_footer(s, 2, TOTAL)

# ========================================================================
# Slide 3 — RAPPEL RÉGLEMENTAIRE PDP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Rappel réglementaire — Plan de Prévention")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "Code du travail — articles R4511-1 et suivants  |  Décret n° 92-158 du 20 février 1992",
             size=14, bold=True, color=VERT)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(3.5), [
    "PDP obligatoire dès qu'une entreprise extérieure (EE) intervient chez une entreprise utilisatrice (EU).",
    "Seuils : intervention ≥ 400 heures sur 12 mois  OU  travaux dangereux (arrêté du 19 mars 1993).",
    "Travaux dangereux : électricité HT/BT, travail en hauteur, soudure, espaces confinés, amiante, levage…",
    "Le PDP est écrit, daté, signé par l'EU et chaque EE, conservé pendant 5 ans minimum.",
    "Mis à jour à chaque modification du périmètre, des risques ou des intervenants.",
], size=17)
# Encadré sanctions
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.7),
                        Inches(12.1), Inches(1.1))
sh.line.color.rgb = JAUNE
sh.line.width = Pt(2)
sh.fill.solid()
sh.fill.fore_color.rgb = RGBColor(0xFF, 0xF7, 0xE0)
tf = sh.text_frame
tf.word_wrap = True
tf.margin_left = Inches(0.2)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
r1 = p.add_run()
r1.text = "Sanctions  "
r1.font.name = "Calibri"; r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = VERT
r2 = p.add_run()
r2.text = ("Absence de PDP = contravention 5e classe (1 500 € / salarié exposé). "
           "En cas d'accident : faute inexcusable de l'employeur, responsabilité pénale du dirigeant.")
r2.font.name = "Calibri"; r2.font.size = Pt(14); r2.font.color.rgb = NOIR
add_footer(s, 3, TOTAL)

# ========================================================================
# Slide 4 — RÉGLEMENTAIRE ICP & PERMIS DE FEU
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Rappel réglementaire — ICP et Permis de feu")
# Colonne ICP
add_text_box(s, Inches(0.6), Inches(1.2), Inches(6), Inches(0.5),
             "Inspection Commune Préalable (ICP)", size=18, bold=True, color=VERT)
add_bullets(s, Inches(0.6), Inches(1.8), Inches(6.1), Inches(4.5), [
    "Obligatoire AVANT toute rédaction du PDP.",
    "Visite sur site, conjointe EU + EE.",
    "Objet : reconnaître les lieux, identifier les risques d'interférences.",
    "Tracée par un compte-rendu signé : c'est l'ICP de l'application.",
    "Pas d'ICP = pas de PDP valable.",
], size=16)
# Colonne Permis de feu
add_text_box(s, Inches(6.9), Inches(1.2), Inches(6), Inches(0.5),
             "Permis de feu", size=18, bold=True, color=VERT)
add_bullets(s, Inches(6.9), Inches(1.8), Inches(6.1), Inches(4.5), [
    "Obligatoire pour tout point chaud : soudure, meulage, découpe, brasage…",
    "Un permis par opération, par jour, par zone.",
    "Signé : Représentant EU + Représentant EE + Exécutant.",
    "Ronde de surveillance imposée 2 h après la fin des travaux.",
    "Émargement de chaque ronde tracé dans l'app.",
], size=16)
# Bande rappel
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.4),
                        Inches(12.1), Inches(0.55))
sh.line.fill.background()
sh.fill.solid()
sh.fill.fore_color.rgb = JAUNE
tf = sh.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "Plan signé = preuve juridique opposable en cas d'accident",
         size=15, bold=True, color=VERT, align=PP_ALIGN.CENTER)
add_footer(s, 4, TOTAL)

# ========================================================================
# Slide 5 — WORKFLOW
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Vue d'ensemble du workflow")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "De la demande client à l'archivage : un parcours unique, sans rupture.",
             size=15, color=GRIS_FONCE)

steps = [
    ("1. Sellsy", "Recherche\npar n° OPP", VERT),
    ("2. ICP", "Visite préalable\n+ photos + GPS", VERT),
    ("3. Risques", "Analyse conjointe\nEU / EE", VERT),
    ("4. PDP", "Héritage auto\ndepuis l'ICP", VERT),
    ("5. Signatures", "Tactile multi-\nsignataires", JAUNE),
    ("6. SharePoint", "Archivage auto\nPDF + DOCX + ZIP", VERT),
]
n = len(steps)
total_w = Inches(12.2)
box_w = Inches(1.75)
gap = (total_w - box_w * n) / (n - 1)
top = Inches(2.8)
left = Inches(0.55)
for i, (titre, corps, col) in enumerate(steps):
    x = left + (box_w + gap) * i
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, Inches(2))
    sh.line.color.rgb = col
    sh.line.width = Pt(2)
    sh.fill.solid()
    sh.fill.fore_color.rgb = BLANC if col == VERT else RGBColor(0xFF, 0xF7, 0xE0)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = titre
    r1.font.name = "Calibri"; r1.font.size = Pt(14); r1.font.bold = True; r1.font.color.rgb = col
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = corps
    r2.font.name = "Calibri"; r2.font.size = Pt(11); r2.font.color.rgb = NOIR
    # flèche entre
    if i < n - 1:
        ax = x + box_w + Inches(0.02)
        aw = gap - Inches(0.04)
        ay = top + Inches(0.9)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, aw, Inches(0.2))
        arr.line.fill.background()
        arr.fill.solid()
        arr.fill.fore_color.rgb = VERT

# 8 statuts en bas
add_text_box(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
             "Cycle de vie d'un PDP — 8 statuts :", size=14, bold=True, color=VERT)
statuts = "Brouillon → Visite préalable → Analyse risque → En signature → Actif → Modifié → Clôturé → Archivé"
add_text_box(s, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
             statuts, size=14, color=NOIR)
add_footer(s, 5, TOTAL)

# ========================================================================
# Slide 6 — LES 6 MODULES
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Les 6 modules de l'application")
mods = [
    ("Tableau de bord", "KPIs en direct : nb PDP actifs, échéances, alertes"),
    ("PDP", "Liste, filtres, recherche, exports PDF / DOCX / ZIP"),
    ("ICP", "Inspections préalables — Mode Express disponible"),
    ("Entreprises", "Répertoire EU / EE, multi-contacts, FDS rattachées"),
    ("Paramètres", "Compte, sync SharePoint, Sellsy, signatures par défaut"),
    ("Aide", "13 tutos consultables — modification admin"),
]
cols = 3
rows = 2
gx = Inches(0.6)
gy = Inches(1.3)
bw = Inches(4.0)
bh = Inches(2.7)
hgap = Inches(0.1)
vgap = Inches(0.15)
for i, (titre, corps) in enumerate(mods):
    r = i // cols
    c = i % cols
    x = gx + (bw + hgap) * c
    y = gy + (bh + vgap) * r
    add_box(s, x, y, bw, bh, titre, corps,
            title_size=18, body_size=14)
add_footer(s, 6, TOTAL)

# ========================================================================
# Slide 7 — ÉTAPE 1 : CRÉER UNE ICP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 1 — Créer une ICP")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.5), [
    "Nouveau → ICP (ou Mode Express).",
    "Taper le n° OPP du dossier Sellsy.",
    "Clic sur l'icône loupe : récupération auto du client, adresse, GPS, contact.",
    "Date de visite, participants EU et EE.",
    "Photos chantier directement depuis l'appareil photo.",
    "Géocodage automatique de l'adresse vers coordonnées GPS.",
], size=16)
add_screenshot_placeholder(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3),
                           "vue ICP — saisie d'un n° OPP avec la recherche Sellsy en cours")
register_screenshot(7, "Vue ICP — saisie d'un nouveau ICP avec la recherche Sellsy ouverte sur un n° OPP, opportunité affichée")
add_footer(s, 7, TOTAL)

# ========================================================================
# Slide 8 — ÉTAPE 2 : RISQUES
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 2 — Identifier les risques d'interférences")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "Catalogue de risques par catégories — sélection en cases à cocher",
             size=15, color=GRIS_FONCE)
# Tableau de catégories
cats = [
    ("Électrique", "Contact direct / indirect, arc électrique, HT/BT"),
    ("Chute", "Hauteur, plain-pied, fouille, toiture"),
    ("Thermique", "Brûlure, point chaud, soudure, gel"),
    ("Chimique", "Produits dangereux, FDS, ATEX"),
    ("Machine / outil", "Coupure, écrasement, projection, bruit"),
    ("Co-activité", "Croisement d'engins, charges suspendues"),
]
top = Inches(1.9)
bw = Inches(6.1)
bh = Inches(1.4)
for i, (n_, d_) in enumerate(cats):
    r = i // 2
    c = i % 2
    x = Inches(0.6) + (bw + Inches(0.15)) * c
    y = top + (bh + Inches(0.12)) * r
    add_box(s, x, y, bw, bh, n_, d_, title_size=15, body_size=13)
# Bande info
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(6.55),
                        Inches(12.1), Inches(0.5))
sh.line.fill.background()
sh.fill.solid(); sh.fill.fore_color.rgb = VERT
tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "Pour chaque risque : niveau 0/1/2  •  mesures de prévention  •  charge EU ou EE",
         size=14, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_footer(s, 8, TOTAL)
register_screenshot(8, "Vue ICP/PDP — onglet Risques avec plusieurs catégories cochées et mesures de prévention saisies")

# ========================================================================
# Slide 9 — ÉTAPE 3 : SIGNATURES ICP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 3 — Signatures ICP")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.5), [
    "Canvas tactile natif : doigt sur smartphone, stylet sur tablette, souris sur PC.",
    "Représentant EU obligatoire.",
    "Un ou plusieurs Représentants EE selon le nombre d'entreprises présentes.",
    "Nom, fonction, date horodatée enregistrés automatiquement.",
    "Bouton Effacer pour recommencer, Valider pour figer la signature.",
    "Signature stockée en base64 — repartie dans le PDP et les exports.",
], size=16)
add_screenshot_placeholder(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3),
                           "écran de signature ICP — canvas tactile avec une signature en cours")
register_screenshot(9, "Vue ICP — bloc Signatures, canvas ouvert avec une signature en cours, nom et fonction renseignés")
add_footer(s, 9, TOTAL)

# ========================================================================
# Slide 10 — ICP -> PDP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 4 — Du ICP au PDP en un clic")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "Tout ce qui a été saisi dans l'ICP est repris automatiquement.",
             size=15, color=GRIS_FONCE)
# Schéma : ICP -> bouton -> PDP
left = Inches(0.6); top = Inches(2.2); w = Inches(3.8); h = Inches(2.5)
add_box(s, left, top, w, h, "ICP — Source",
        "• Client, adresse, GPS\n• Participants EU / EE\n• Risques identifiés\n• Photos\n• Signatures",
        title_size=18, body_size=13)
# flèche
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.6), Inches(3.1),
                         Inches(1.2), Inches(0.7))
arr.line.fill.background(); arr.fill.solid(); arr.fill.fore_color.rgb = JAUNE
add_text_box(s, Inches(4.5), Inches(3.85), Inches(1.4), Inches(0.5),
             "Générer\nPDP", size=12, bold=True, color=VERT, align=PP_ALIGN.CENTER)
add_box(s, Inches(6.0), top, w, h, "PDP — Hérité",
        "• Données pré-remplies\n• Risques recopiés\n• Mesures à compléter\n• Permis de feu si besoin\n• Annexes / FDS",
        title_size=18, body_size=13, border=JAUNE, title_color=JAUNE)
add_screenshot_placeholder(s, Inches(10.2), Inches(2.2), Inches(2.7), Inches(2.5),
                           "bouton Générer le PDP depuis un ICP signé")
register_screenshot(10, "Vue ICP signé — bouton 'Générer le PDP' visible avec confirmation modale")
# bandeau gain
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.3),
                        Inches(12.1), Inches(1.2))
sh.line.fill.background(); sh.fill.solid(); sh.fill.fore_color.rgb = VERT
tf = sh.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.3)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r1 = p.add_run(); r1.text = "Zéro double-saisie\n"
r1.font.name = "Calibri"; r1.font.size = Pt(22); r1.font.bold = True; r1.font.color.rgb = JAUNE
r2 = p.add_run(); r2.text = "Gain de temps moyen estimé : 20 à 30 minutes par dossier"
r2.font.name = "Calibri"; r2.font.size = Pt(16); r2.font.color.rgb = BLANC
add_footer(s, 10, TOTAL)

# ========================================================================
# Slide 11 — ÉTAPE 5 : COMPLÉTER LE PDP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 5 — Compléter le PDP")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "Une fiche PDP organisée en 10 onglets — navigation libre, sauvegarde automatique.",
             size=15, color=GRIS_FONCE)
onglets = [
    ("Général", "OPP, client, dates, lieu"),
    ("Entreprises", "EU + EE, contacts, effectifs"),
    ("Visite", "Compte-rendu ICP, photos"),
    ("Risques", "Catalogue, niveaux, mesures"),
    ("Urgence", "Consignes, numéros, points de rassemblement"),
    ("Signatures", "EU + chaque EE"),
    ("Permis feu", "1 à N permis liés au PDP"),
    ("Émargement", "Rondes de surveillance"),
    ("Pictogrammes", "EPI / interdiction / obligation"),
    ("Annexes", "FDS, modes opératoires, plans"),
]
cols = 5
rows = 2
bw = Inches(2.42)
bh = Inches(1.5)
top = Inches(1.9)
left = Inches(0.6)
hgap = Inches(0.08)
vgap = Inches(0.12)
for i, (n_, d_) in enumerate(onglets):
    r = i // cols
    c = i % cols
    x = left + (bw + hgap) * c
    y = top + (bh + vgap) * r
    add_box(s, x, y, bw, bh, n_, d_, title_size=13, body_size=11)
add_screenshot_placeholder(s, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.6),
                           "vue PDP ouvert avec la barre des 10 onglets visible en haut")
register_screenshot(11, "Vue PDP — fiche ouverte montrant la barre des 10 onglets en haut, focus sur l'onglet Général ou Risques")
add_footer(s, 11, TOTAL)

# ========================================================================
# Slide 12 — PERMIS DE FEU
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Permis de feu et rondes de surveillance")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.5), [
    "Création depuis l'onglet Permis feu du PDP.",
    "Date, horaires, zone d'intervention, type de point chaud.",
    "3 signatures obligatoires : Rep EU + Rep EE + Exécutant.",
    "Émargement par ronde de surveillance (jusqu'à 2 h après).",
    "Chaque ronde signée, horodatée, géolocalisée.",
    "Traçabilité complète exportée dans le PDP.",
], size=16)
add_screenshot_placeholder(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3),
                           "onglet Permis de feu avec 3 signatures + tableau d'émargement des rondes")
register_screenshot(12, "Vue PDP — onglet Permis de feu, avec un permis créé, ses 3 signatures et le tableau d'émargement des rondes")
add_footer(s, 12, TOTAL)

# ========================================================================
# Slide 13 — BIBLIOTHÈQUE FDS
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Bibliothèque FDS centralisée")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.5), [
    "Une seule bibliothèque pour tous les dossiers : plus de fichiers éparpillés.",
    "Marquage 'Standard' : la FDS s'attache automatiquement à chaque nouveau PDP.",
    "FDS personnalisées : à rattacher manuellement à un dossier spécifique.",
    "Rangement par entreprise et par produit, recherche instantanée.",
], size=16)
# 3 cartes
top = Inches(4.0); bh = Inches(2.6)
add_box(s, Inches(0.6), top, Inches(4.0), bh,
        "FDS Standard",
        "Produits utilisés en permanence par AGRIWATT.\nAuto-rattachées à TOUS les PDP.",
        title_size=16, body_size=13, border=JAUNE, title_color=JAUNE)
add_box(s, Inches(4.75), top, Inches(4.0), bh,
        "FDS spécifiques EE",
        "Produits apportés par les entreprises extérieures.\nRattachées au PDP concerné.",
        title_size=16, body_size=13)
add_box(s, Inches(8.9), top, Inches(4.0), bh,
        "Mise à jour",
        "Pictogrammes, dangers, mesures à jour selon les FDS fournisseur.\nRéforme CLP intégrée.",
        title_size=16, body_size=13)
add_footer(s, 13, TOTAL)

# ========================================================================
# Slide 14 — SELLSY : AVANT / APRÈS
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Intégration Sellsy — Avant / Après")
# Avant
add_text_box(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(0.5),
             "AVANT — saisie manuelle", size=18, bold=True, color=GRIS_FONCE)
add_bullets(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.5), [
    "Ouvrir Sellsy dans un autre onglet.",
    "Chercher le client, copier l'adresse.",
    "Coller, vérifier, corriger les fautes.",
    "Trouver les coordonnées du contact.",
    "Recopier téléphone + email.",
    "~ 8 à 10 minutes par dossier.",
], size=14, bullet_color=GRIS_MOYEN)
# Après
add_text_box(s, Inches(6.9), Inches(1.2), Inches(6.0), Inches(0.5),
             "APRÈS — recherche Sellsy intégrée", size=18, bold=True, color=VERT)
add_bullets(s, Inches(6.9), Inches(1.8), Inches(6.1), Inches(4.5), [
    "Taper le n° OPP dans l'app.",
    "Cliquer sur la loupe.",
    "Sélectionner l'opportunité affichée.",
    "Tout est pré-rempli : adresse, GPS, contact, téléphone, email.",
    "L'entreprise est créée automatiquement dans le répertoire local.",
    "~ 30 secondes.",
], size=14, bullet_color=JAUNE)
add_screenshot_placeholder(s, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.7),
                           "résultats Sellsy affichés après recherche par OPP — opportunité prête à être sélectionnée")
register_screenshot(14, "Vue ICP — résultats de la recherche Sellsy après saisie d'un n° OPP, liste d'opportunités affichée")
add_footer(s, 14, TOTAL)

# ========================================================================
# Slide 15 — MOBILITÉ TERRAIN
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Mobilité terrain")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "L'app fonctionne là où vous travaillez — même sans réseau.",
             size=15, color=GRIS_FONCE)
items = [
    ("100 % offline", "Service Worker + IndexedDB : tout est utilisable sans connexion."),
    ("Photos natives", "Appareil photo du smartphone, compression auto, intégrées au dossier."),
    ("GPS automatique", "Coordonnées du chantier capturées en 1 clic."),
    ("Sync multi-appareils", "Sync SharePoint dès le retour wifi/4G, 2 sec après chaque enregistrement."),
]
top = Inches(2.0)
bw = Inches(6.0)
bh = Inches(1.8)
hgap = Inches(0.15)
vgap = Inches(0.2)
for i, (n_, d_) in enumerate(items):
    r = i // 2; c = i % 2
    x = Inches(0.6) + (bw + hgap) * c
    y = top + (bh + vgap) * r
    add_box(s, x, y, bw, bh, n_, d_, title_size=18, body_size=14)
add_screenshot_placeholder(s, Inches(0.6), Inches(5.95), Inches(12.1), Inches(1.0),
                           "vue smartphone — appli installée en PWA, mode hors-ligne actif")
register_screenshot(15, "Smartphone — appli installée en PWA, icône offline visible, vue PDP ouverte avec photo chantier")
add_footer(s, 15, TOTAL)

# ========================================================================
# Slide 16 — EXPORT & ARCHIVAGE
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Export et archivage")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(4.0), [
    "Export PDF — pour signature papier ou envoi client.",
    "Export Word (DOCX) — pour modifications avant validation.",
    "Export ZIP — dossier complet : PDP + annexes + FDS + photos.",
    "Archivage automatique SharePoint à la clôture.",
    "Naming standardisé pour la recherche ultérieure.",
], size=16)
# Bloc naming
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.4),
                        Inches(5.8), Inches(2.0))
sh.line.color.rgb = VERT; sh.line.width = Pt(1.5)
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF2, 0xF8, 0xF4)
tf = sh.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
p1 = tf.paragraphs[0]
r1 = p1.add_run(); r1.text = "Naming d'export\n"
r1.font.name = "Calibri"; r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = VERT
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "OPP_PDP-YYYY-MM-XXXX_titre"
r2.font.name = "Consolas"; r2.font.size = Pt(15); r2.font.bold = True; r2.font.color.rgb = NOIR
p3 = tf.add_paragraph()
r3 = p3.add_run(); r3.text = "\nExemple :\nOPP12345_PDP-2026-06-0042_Toiture-Mairie.zip"
r3.font.name = "Consolas"; r3.font.size = Pt(12); r3.font.color.rgb = GRIS_FONCE
# Chemins SharePoint
add_text_box(s, Inches(7.0), Inches(3.6), Inches(5.8), Inches(0.4),
             "Chemins d'archivage :", size=14, bold=True, color=VERT)
add_text_box(s, Inches(7.0), Inches(4.0), Inches(5.8), Inches(0.4),
             "SuiviPDP/PDP/Exports/", size=13, font="Consolas")
add_text_box(s, Inches(7.0), Inches(4.45), Inches(5.8), Inches(0.4),
             "SuiviPDP/ICP/Exports/", size=13, font="Consolas")
add_screenshot_placeholder(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.5),
                           "boîte de dialogue d'export proposant PDF / DOCX / ZIP")
register_screenshot(16, "Vue PDP — boîte de dialogue d'export avec les 3 boutons PDF, DOCX, ZIP visibles")
add_footer(s, 16, TOTAL)

# ========================================================================
# Slide 17 — BONNES PRATIQUES
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Bonnes pratiques — Do's & Don'ts")
# DO
add_text_box(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(0.5),
             "À FAIRE", size=20, bold=True, color=VERT)
add_bullets(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.5), [
    "Synchroniser en wifi AVANT de partir sur chantier.",
    "Prendre les photos systématiquement à chaque visite.",
    "Faire signer EU et EE AVANT l'accès chantier.",
    "Vérifier que les FDS sont à jour avant le démarrage.",
    "Sauvegarder régulièrement — l'auto-save existe, mais un coup de pouce ne nuit pas.",
], size=15, bullet_color=VERT)
# DON'T
add_text_box(s, Inches(6.9), Inches(1.2), Inches(6.0), Inches(0.5),
             "À ÉVITER", size=20, bold=True, color=RGBColor(0xB0, 0x00, 0x00))
add_bullets(s, Inches(6.9), Inches(1.8), Inches(6.1), Inches(4.5), [
    "Démarrer un chantier sans PDP signé.",
    "Skipper l'ICP : pas d'ICP = PDP fragilisé juridiquement.",
    "Lancer un point chaud sans permis de feu.",
    "Oublier les rondes de surveillance après soudure / meulage.",
    "Modifier un PDP clôturé : créer une nouvelle version.",
], size=15, bullet_color=RGBColor(0xB0, 0x00, 0x00))
add_footer(s, 17, TOTAL)

# ========================================================================
# Slide 18 — AIDE INTÉGRÉE
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Aide intégrée — 13 tutoriels")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.5), [
    "Module Aide accessible depuis le menu principal.",
    "13 tutoriels couvrant tous les cas d'usage métier.",
    "Recherche par mot-clé dans les tutoriels.",
    "Pas-à-pas illustrés, navigation libre.",
    "Modification possible par administrateur (nouveauté).",
    "Mise à jour continue selon les retours terrain.",
], size=16)
add_screenshot_placeholder(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3),
                           "vue Aide — liste des 13 tutoriels avec barre de recherche")
register_screenshot(18, "Vue Aide — liste des 13 tutoriels affichée avec barre de recherche en haut")
add_footer(s, 18, TOTAL)

# ========================================================================
# Slide 19 — Q&A
# ========================================================================
s = prs.slides.add_slide(BLANK)
# fond vert
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background()
bg.fill.solid(); bg.fill.fore_color.rgb = VERT
# bande jaune
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.0), SW, Inches(0.08))
b.line.fill.background(); b.fill.solid(); b.fill.fore_color.rgb = JAUNE
add_text_box(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.5),
             "Vos questions ?", size=72, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.8),
             "On en discute maintenant — ou à tout moment après la formation.",
             size=20, color=JAUNE, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
             "T. HOCEDEZ", size=26, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
             "Responsable QSE — AGRIWATT", size=20, color=BLANC, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
             "Merci pour votre attention.", size=18, color=JAUNE, align=PP_ALIGN.CENTER, font="Calibri")

# ========================================================================
# Slide 20 — ANNEXE
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Annexe — Raccourcis et cas particuliers")
items = [
    ("Mode Express ICP",
     "Pour les petits chantiers : ICP simplifié, saisie réduite, génération PDP directe."),
    ("Multi-EE",
     "Plusieurs entreprises extérieures sur le même PDP : ajout dans l'onglet Entreprises, signature par EE."),
    ("FDS personnalisée",
     "Ajouter une FDS non standard : Bibliothèque FDS → Nouvelle → décocher 'Standard' → rattacher au PDP."),
    ("Modification post-clôture",
     "Statut 'Modifié' : crée une nouvelle version du PDP avec historique, sans écraser la précédente."),
    ("Sync forcée",
     "Paramètres → Sync SharePoint → Lancer la sync manuellement après une longue absence wifi."),
    ("Re-géocodage",
     "Si l'adresse change : ICP/PDP → Re-géocoder → mise à jour des coordonnées GPS."),
]
top = Inches(1.2)
bw = Inches(6.05)
bh = Inches(1.75)
hgap = Inches(0.15)
vgap = Inches(0.12)
for i, (n_, d_) in enumerate(items):
    r = i // 2; c = i % 2
    x = Inches(0.6) + (bw + hgap) * c
    y = top + (bh + vgap) * r
    add_box(s, x, y, bw, bh, n_, d_, title_size=15, body_size=12)
add_footer(s, 20, TOTAL)

# Sauvegarde
out = r"C:\Users\ThibaultHOCEDEZ\Documents\suivipdp\Formation_SuiviPDP_2026-06-01.pptx"
prs.save(out)
print(f"OK -> {out}")
print(f"Slides : {len(prs.slides)}")
print(f"Placeholders screenshots : {len(SCREENSHOTS)}")

# Fichier liste screenshots
txt_out = r"C:\Users\ThibaultHOCEDEZ\Documents\suivipdp\Liste_screenshots_a_prendre.txt"
with open(txt_out, "w", encoding="utf-8") as f:
    f.write("Liste des captures d'écran à prendre — Formation Suivi PDP/ICP\n")
    f.write("Présentation du 01/06/2026 — T. HOCEDEZ\n")
    f.write("=" * 70 + "\n\n")
    f.write("Application : C:\\Users\\ThibaultHOCEDEZ\\Documents\\suivipdp\\index.html\n\n")
    f.write("Conseils :\n")
    f.write("- Format conseillé : PNG, largeur 1400 à 1800 px.\n")
    f.write("- Masquer les données client réelles si captures partagées hors AGRIWATT.\n")
    f.write("- Une fois prises, glisser-déposer la capture dans le rectangle gris correspondant.\n\n")
    f.write("-" * 70 + "\n\n")
    for num, desc in SCREENSHOTS:
        f.write(f"Slide {num:>2}  :  {desc}\n\n")
    f.write("-" * 70 + "\n")
    f.write(f"Total : {len(SCREENSHOTS)} captures à prendre.\n")
print(f"Liste captures : {txt_out}")
