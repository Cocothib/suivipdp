"""Génération du PowerPoint de formation — Module ICP (équipes Maintenance) - AGRIWATT
Version centrée sur l'Inspection Commune Préalable pour les visites préventives de maintenance.
Aucune mention du PDP. Met en avant Mode Express (spécial maintenance) vs ICP complète.
Date présentation : 11/06/2026
Auteur : T. HOCEDEZ
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image as PILImage

CAP_DIR = r"C:\Users\ThibaultHOCEDEZ\Documents\suivipdp\captures_formation"

# Couleurs Agriwatt
VERT = RGBColor(0x00, 0x69, 0x3E)
JAUNE = RGBColor(0xFF, 0xB8, 0x1C)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
GRIS_CLAIR = RGBColor(0xEE, 0xEE, 0xEE)
GRIS_MOYEN = RGBColor(0xBB, 0xBB, 0xBB)
GRIS_FONCE = RGBColor(0x55, 0x55, 0x55)
NOIR = RGBColor(0x22, 0x22, 0x22)
ROUGE = RGBColor(0xB0, 0x00, 0x00)

DATE_PRESENTATION = "11/06/2026"
DATE_LONGUE = "11 juin 2026"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]
TOTAL = 15


def set_text(tf, text, size=18, bold=False, color=NOIR, align=PP_ALIGN.LEFT, font="Calibri"):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def add_text_box(slide, left, top, width, height, text, size=18, bold=False,
                 color=NOIR, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    set_text(tf, text, size=size, bold=bold, color=color, align=align, font=font)
    return tb


def add_bullets(slide, left, top, width, height, items, size=18, color=NOIR,
                bullet_color=VERT, line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        rb = p.add_run()
        rb.text = "▸  "
        rb.font.name = "Calibri"
        rb.font.size = Pt(size)
        rb.font.bold = True
        rb.font.color.rgb = bullet_color
        rt = p.add_run()
        rt.text = item
        rt.font.name = "Calibri"
        rt.font.size = Pt(size)
        rt.font.color.rgb = color
    return tb


def add_header(slide, title):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.9))
    band.line.fill.background()
    band.fill.solid()
    band.fill.fore_color.rgb = VERT
    add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.6),
                 title, size=24, bold=True, color=BLANC, anchor=MSO_ANCHOR.MIDDLE)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.9), SW, Inches(0.05))
    accent.line.fill.background()
    accent.fill.solid()
    accent.fill.fore_color.rgb = JAUNE


def add_footer(slide, idx, total):
    add_text_box(slide, Inches(0.5), Inches(7.1), Inches(9), Inches(0.3),
                 f"Suivi PDP — Module ICP — Maintenance — {DATE_PRESENTATION}",
                 size=10, color=GRIS_FONCE)
    add_text_box(slide, Inches(11.3), Inches(7.1), Inches(1.5), Inches(0.3),
                 f"{idx} / {total}", size=10, color=GRIS_FONCE, align=PP_ALIGN.RIGHT)


def add_screenshot_placeholder(slide, left, top, width, height, description):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    sh.line.color.rgb = GRIS_MOYEN
    sh.line.width = Pt(1.5)
    sh.fill.solid()
    sh.fill.fore_color.rgb = GRIS_CLAIR
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    set_text(tf, f"[ Capture à insérer ]\n{description}",
             size=14, bold=True, color=GRIS_FONCE, align=PP_ALIGN.CENTER)


def add_box(slide, left, top, width, height, title, body, fill=BLANC,
            border=VERT, title_color=VERT, body_color=NOIR, title_size=16, body_size=13):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    sh.line.color.rgb = border
    sh.line.width = Pt(1.5)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = "Calibri"
    r1.font.size = Pt(title_size)
    r1.font.bold = True
    r1.font.color.rgb = title_color
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    p2.line_spacing = 1.1
    r2 = p2.add_run()
    r2.text = body
    r2.font.name = "Calibri"
    r2.font.size = Pt(body_size)
    r2.font.color.rgb = body_color


def band_bottom(slide, text, top=Inches(6.4), fill=JAUNE, color=VERT, size=15):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(12.1), Inches(0.55))
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    tf = sh.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(tf, text, size=size, bold=True, color=color, align=PP_ALIGN.CENTER)


def add_image_fit(slide, left, top, width, height, filename, desc=""):
    """Insère la capture en conservant le ratio, centrée dans la zone, avec un fin cadre gris.
    Retombe sur le placeholder gris si le fichier est absent."""
    path = os.path.join(CAP_DIR, filename) if filename else None
    if not (path and os.path.exists(path)):
        add_screenshot_placeholder(slide, left, top, width, height, desc)
        return
    iw, ih = PILImage.open(path).size
    ar = iw / ih
    box_ar = width / height
    if ar > box_ar:
        w = int(width); h = int(width / ar)
    else:
        h = int(height); w = int(height * ar)
    x = int(left) + (int(width) - w) // 2
    y = int(top) + (int(height) - h) // 2
    pad = Emu(9525)  # ~0.01"
    frame = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x - pad, y - pad, w + 2 * pad, h + 2 * pad)
    frame.line.color.rgb = GRIS_MOYEN; frame.line.width = Pt(1)
    frame.fill.background()
    slide.shapes.add_picture(path, x, y, width=w, height=h)


SCREENSHOTS = []
def register_screenshot(slide_num, desc):
    SCREENSHOTS.append((slide_num, desc))


# ========================================================================
# Slide 1 — COUVERTURE
# ========================================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background(); bg.fill.solid(); bg.fill.fore_color.rgb = VERT
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(5.2), SW, Inches(0.08))
b.line.fill.background(); b.fill.solid(); b.fill.fore_color.rgb = JAUNE
logo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.5), Inches(2.4), Inches(0.9))
logo.line.color.rgb = JAUNE; logo.line.width = Pt(1.5)
logo.fill.solid(); logo.fill.fore_color.rgb = VERT
tf = logo.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, "AGRIWATT", size=20, bold=True, color=JAUNE, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.6), Inches(2.2), Inches(12), Inches(1.2),
             "Suivi PDP — Module ICP", size=50, bold=True, color=BLANC)
add_text_box(s, Inches(0.6), Inches(3.3), Inches(12), Inches(0.8),
             "L'Inspection Commune Préalable au service de vos visites préventives",
             size=24, color=JAUNE)
add_text_box(s, Inches(0.6), Inches(4.2), Inches(12), Inches(0.7),
             "Formation équipes Maintenance", size=22, color=BLANC)
add_text_box(s, Inches(0.6), Inches(5.5), Inches(12), Inches(0.6),
             DATE_LONGUE, size=20, bold=True, color=BLANC)
add_text_box(s, Inches(0.6), Inches(6.1), Inches(12), Inches(0.5),
             "T. HOCEDEZ — Responsable QSE", size=18, color=BLANC)

# ========================================================================
# Slide 2 — POURQUOI CET OUTIL POUR VOS VISITES
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Pourquoi cet outil pour vos visites préventives ?")
box_w = Inches(5.8); box_h = Inches(2.5); gap = Inches(0.3)
left1 = Inches(0.6); left2 = Inches(0.6) + box_w + gap
top1 = Inches(1.3); top2 = top1 + box_h + gap
add_box(s, left1, top1, box_w, box_h, "Simple et rapide",
        "Mode Express dédié aux visites simples : seuls les champs utiles s'affichent. "
        "Une inspection tracée en quelques minutes, depuis le terrain.", title_size=20, body_size=15)
add_box(s, left2, top1, box_w, box_h, "Conformité",
        "Inspection commune des lieux conforme au Code du travail (art. R4512-2 et suivants). "
        "Compte-rendu signé = preuve de la reconnaissance des risques.", title_size=20, body_size=15)
add_box(s, left1, top2, box_w, box_h, "Mobilité terrain",
        "Application installée sur smartphone / tablette. 100 % hors-ligne, photos natives, "
        "GPS automatique, signature tactile sur site.", title_size=20, body_size=15)
add_box(s, left2, top2, box_w, box_h, "Traçabilité",
        "Historique complet, archivage automatique à la validation, export PDF / Word, "
        "synchronisation SharePoint partagée.", title_size=20, body_size=15)
add_footer(s, 2, TOTAL)

# ========================================================================
# Slide 3 — QU'EST-CE QU'UNE ICP ?
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Qu'est-ce qu'une ICP ?")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "Inspection Commune Préalable — Code du travail, art. R4512-2 à R4512-5",
             size=15, bold=True, color=VERT)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(12), Inches(3.6), [
    "Visite des lieux réalisée AVANT l'intervention, en commun avec le site (entreprise utilisatrice).",
    "Objet : reconnaître les lieux, délimiter la zone, identifier les risques d'interférences (co-activité).",
    "Pour la maintenance préventive : un point d'entrée systématique avant chaque visite à risque.",
    "On définit ensemble : accès, circulation, consignes de sécurité, mesures de prévention.",
    "Tracée par un compte-rendu daté et signé — c'est exactement l'ICP de l'application.",
], size=17)
band_bottom(s, "Reconnaître les lieux et les risques AVANT d'intervenir — et le prouver",
            top=Inches(5.9), fill=JAUNE, color=VERT)
add_footer(s, 3, TOTAL)

# ========================================================================
# Slide 4 — LE PARCOURS D'UNE ICP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Le parcours d'une ICP")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "Cinq étapes, un seul outil, sans rupture — du terrain à l'archivage.",
             size=15, color=GRIS_FONCE)
steps = [
    ("1. Créer", "Sellsy (OPP)\nou manuel", VERT),
    ("2. Reconnaître", "Lieux, accès,\nconsignes", VERT),
    ("3. Risques", "Interférences\n+ photos + GPS", VERT),
    ("4. Signatures", "Site + intervenant\n(tactile)", JAUNE),
    ("5. Archivage", "Export PDF/Word\n+ SharePoint", VERT),
]
n = len(steps)
total_w = Inches(12.2); box_w = Inches(2.05)
gap = (total_w - box_w * n) / (n - 1)
top = Inches(2.8); left = Inches(0.55)
for i, (titre, corps, col) in enumerate(steps):
    x = left + (box_w + gap) * i
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, Inches(2))
    sh.line.color.rgb = col; sh.line.width = Pt(2)
    sh.fill.solid()
    sh.fill.fore_color.rgb = BLANC if col == VERT else RGBColor(0xFF, 0xF7, 0xE0)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = titre
    r1.font.name = "Calibri"; r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = col
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = corps
    r2.font.name = "Calibri"; r2.font.size = Pt(11); r2.font.color.rgb = NOIR
    if i < n - 1:
        ax = x + box_w + Inches(0.02); aw = gap - Inches(0.04); ay = top + Inches(0.9)
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay, aw, Inches(0.2))
        arr.line.fill.background(); arr.fill.solid(); arr.fill.fore_color.rgb = VERT
add_text_box(s, Inches(0.6), Inches(5.3), Inches(12), Inches(0.4),
             "Cycle de vie d'une ICP :", size=14, bold=True, color=VERT)
add_text_box(s, Inches(0.6), Inches(5.8), Inches(12), Inches(0.8),
             "Effectuée  →  Validée  →  Archivée automatiquement", size=16, color=NOIR)
add_footer(s, 4, TOTAL)

# ========================================================================
# Slide 5 — DEUX FAÇONS DE DÉMARRER
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Deux façons de démarrer une ICP")
add_text_box(s, Inches(0.6), Inches(1.1), Inches(6.0), Inches(0.45),
             "Avec Sellsy (n° OPP)", size=18, bold=True, color=VERT)
add_bullets(s, Inches(0.6), Inches(1.6), Inches(6.0), Inches(2.0), [
    "Quand la visite est rattachée à un dossier Sellsy.",
    "Taper le n° OPP, cliquer sur la loupe.",
    "Client, adresse, GPS et contact pré-remplis.",
    "Zéro recopie, zéro faute de frappe.",
], size=14, bullet_color=VERT)
add_text_box(s, Inches(0.6), Inches(3.85), Inches(6.0), Inches(0.45),
             "Création manuelle", size=18, bold=True, color=JAUNE)
add_bullets(s, Inches(0.6), Inches(4.35), Inches(6.0), Inches(2.0), [
    "Quand il n'y a pas de dossier Sellsy (visite interne).",
    "Saisir directement le site, l'adresse, le contact.",
    "Géocodage de l'adresse → GPS en 1 clic.",
    "Idéal combiné au Mode Express (slide suivante).",
], size=14, bullet_color=JAUNE)
add_image_fit(s, Inches(6.95), Inches(1.2), Inches(5.85), Inches(5.3),
              "slide14_sellsy_resultats.png",
              "vue ICP — recherche Sellsy par n° OPP")
register_screenshot(5, "Vue ICP — écran de création montrant le champ n° OPP (recherche Sellsy) et la saisie manuelle du site")
add_footer(s, 5, TOTAL)

# ========================================================================
# Slide 6 — MODE EXPRESS vs ICP COMPLÈTE
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Mode Express vs ICP complète")
add_text_box(s, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             "Un simple interrupteur en haut de l'ICP. Choisissez selon la complexité de la visite.",
             size=15, color=GRIS_FONCE)
# Carte Mode Express (mise en avant maintenance)
add_box(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.4),
        "⚡ Mode Express  —  recommandé maintenance",
        "Masque les sections avancées pour aller à l'essentiel :\n\n"
        "• Site, contact, date de visite\n"
        "• Délimitation / circulation / consignes\n"
        "• Risques d'interférences\n"
        "• Photos + GPS\n"
        "• Signatures site + intervenant\n\n"
        "Idéal pour une visite préventive simple, un seul intervenant.",
        title_size=17, body_size=14, border=JAUNE, title_color=JAUNE,
        fill=RGBColor(0xFF, 0xF7, 0xE0))
# Carte ICP complète
add_box(s, Inches(6.9), Inches(1.8), Inches(6.0), Inches(4.4),
        "ICP complète",
        "Affiche en plus les sections de co-activité :\n\n"
        "• Autres intervenants / sous-traitants (EE)\n"
        "• Participants additionnels\n"
        "• Signatures de chaque intervenant supplémentaire\n\n"
        "À utiliser quand plusieurs entreprises interviennent en même temps "
        "sur la même zone.",
        title_size=17, body_size=14)
band_bottom(s, "Le Mode Express se mémorise : activez-le une fois, il reste actif pour vos prochaines visites",
            top=Inches(6.45), fill=VERT, color=BLANC, size=13)
register_screenshot(6, "Vue ICP — interrupteur 'Mode Express' en haut, une fois activé (sections avancées masquées)")
add_footer(s, 6, TOTAL)

# ========================================================================
# Slide 7 — ÉTAPE 1 : CRÉER L'ICP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 1 — Créer l'ICP")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.3), [
    "Menu : Nouveau → ICP.",
    "Activer le Mode Express pour une visite simple.",
    "Renseigner le site et le contact (via Sellsy ou manuellement).",
    "Date de visite et représentant côté site.",
    "Représentant côté intervenant (vous / l'équipe maintenance).",
    "Géocodage de l'adresse vers coordonnées GPS.",
], size=16)
add_image_fit(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3),
              "slide07_icp_nouvelle_recherche_sellsy.png",
              "vue ICP en Mode Express — champs généraux")
register_screenshot(7, "Vue ICP en Mode Express — saisie des informations générales (site, date, représentants)")
add_footer(s, 7, TOTAL)

# ========================================================================
# Slide 8 — ÉTAPE 2 : DÉLIMITATION / CIRCULATION / CONSIGNES
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 2 — Délimitation, circulation & consignes")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.5),
             "Points obligatoires de l'inspection commune (art. R4512-3 & R4512-4) — à cocher sur site.",
             size=15, color=GRIS_FONCE)
pts = [
    ("Secteur délimité", "Zone d'intervention clairement identifiée."),
    ("Zones de danger", "Dangers matérialisés et signalés."),
    ("Circulation", "Voies piétons et engins indiquées."),
    ("Voies d'accès", "Accès au poste de travail définis."),
    ("Consignes sécurité", "Consignes du site communiquées à l'équipe."),
    ("Précisions", "Champ libre pour toute remarque utile."),
]
top = Inches(1.9); bw = Inches(6.05); bh = Inches(1.35)
for i, (n_, d_) in enumerate(pts):
    r = i // 2; c = i % 2
    x = Inches(0.6) + (bw + Inches(0.15)) * c
    y = top + (bh + Inches(0.12)) * r
    add_box(s, x, y, bw, bh, n_, d_, title_size=15, body_size=13)
band_bottom(s, "Chaque point coché trace ce qui a été vérifié et abordé avec le site",
            top=Inches(6.5), fill=VERT, color=BLANC, size=14)
register_screenshot(8, "Vue ICP — section Délimitation/circulation/consignes avec les points cochés et un commentaire")
add_footer(s, 8, TOTAL)

# ========================================================================
# Slide 9 — ÉTAPE 3 : RISQUES D'INTERFÉRENCES
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 3 — Identifier les risques")
add_text_box(s, Inches(0.6), Inches(1.15), Inches(6.1), Inches(0.5),
             "Catalogue de risques par catégories — cases à cocher :",
             size=14, bold=True, color=VERT)
add_bullets(s, Inches(0.6), Inches(1.75), Inches(6.1), Inches(4.3), [
    "Électrique — contact direct/indirect, arc, HT/BT",
    "Chute — hauteur, plain-pied, fouille, toiture",
    "Thermique — brûlure, point chaud, gel",
    "Chimique — produits dangereux, FDS, ATEX",
    "Machine / outil — coupure, écrasement, bruit",
    "Co-activité — engins, charges suspendues",
], size=15, line_spacing=1.3)
add_image_fit(s, Inches(6.95), Inches(1.15), Inches(5.85), Inches(5.1),
              "slide08_risques_mesures.png",
              "vue ICP — onglet Risques")
band_bottom(s, "Pour chaque risque : niveau de gravité  •  mesures de prévention  •  qui agit",
            top=Inches(6.45), fill=VERT, color=BLANC, size=14)
register_screenshot(9, "Vue ICP — onglet Risques avec plusieurs catégories cochées et mesures saisies")
add_footer(s, 9, TOTAL)

# ========================================================================
# Slide 10 — ÉTAPE 4 : PHOTOS & GPS
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 4 — Photos & GPS terrain")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.3), [
    "Photos prises directement avec l'appareil du smartphone / tablette.",
    "Compression automatique, rattachées à l'ICP.",
    "Annoter / légender chaque photo (zone, point de vigilance).",
    "Coordonnées GPS du site capturées en 1 clic.",
    "Tout reste disponible même sans réseau.",
], size=16)
add_image_fit(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3),
              "slide15_mobile_pwa.png",
              "vue ICP sur mobile — photos de la visite + GPS")
register_screenshot(10, "Vue ICP — photos de la visite ajoutées, avec légendes et coordonnées GPS renseignées")
add_footer(s, 10, TOTAL)

# ========================================================================
# Slide 11 — ÉTAPE 5 : SIGNATURES ICP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Étape 5 — Signatures de l'ICP")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(5.3), [
    "Canvas tactile : doigt sur smartphone, stylet sur tablette, souris sur PC.",
    "Représentant du site (entreprise utilisatrice) : obligatoire.",
    "Représentant intervenant (équipe maintenance).",
    "Nom, fonction et date horodatée enregistrés automatiquement.",
    "Boutons Effacer / Valider pour figer la signature.",
    "En Mode Express : seules ces 2 signatures — pas de signataires supplémentaires.",
], size=16)
add_image_fit(s, Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.3),
              "slide09_icp_signatures_canvas.png",
              "écran de signature ICP — canvas tactile")
register_screenshot(11, "Vue ICP — bloc Signatures, canvas ouvert, nom et fonction renseignés")
add_footer(s, 11, TOTAL)

# ========================================================================
# Slide 12 — EXPORT & ARCHIVAGE DE L'ICP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Export et archivage de l'ICP")
add_bullets(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(4.2), [
    "Export PDF — pour diffusion ou signature papier.",
    "Export Word (DOCX) — si une retouche est nécessaire.",
    "Passage en statut Validée = archivage automatique.",
    "Dépôt SharePoint partagé, accessible à toute l'équipe.",
    "Nom de fichier standardisé pour la recherche ultérieure.",
], size=16)
sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.4), Inches(5.8), Inches(2.0))
sh.line.color.rgb = VERT; sh.line.width = Pt(1.5)
sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0xF2, 0xF8, 0xF4)
tf = sh.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.1)
p1 = tf.paragraphs[0]
r1 = p1.add_run(); r1.text = "Archivage automatique\n"
r1.font.name = "Calibri"; r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = VERT
p2 = tf.add_paragraph()
r2 = p2.add_run(); r2.text = "SuiviPDP/ICP/Exports/"
r2.font.name = "Consolas"; r2.font.size = Pt(14); r2.font.bold = True; r2.font.color.rgb = NOIR
p3 = tf.add_paragraph()
r3 = p3.add_run(); r3.text = "\nPDF + Word déposés dès la validation."
r3.font.name = "Calibri"; r3.font.size = Pt(12); r3.font.color.rgb = GRIS_FONCE
add_screenshot_placeholder(s, Inches(0.6), Inches(5.4), Inches(12.1), Inches(1.5),
                           "vue ICP — boutons d'export PDF / Word")
register_screenshot(12, "Vue ICP — boutons d'export PDF et Word visibles sur une ICP validée")
add_footer(s, 12, TOTAL)

# ========================================================================
# Slide 13 — MOBILITÉ, AIDE & AMÉLIORATION CONTINUE
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Mobilité, aide & amélioration continue")
items = [
    ("100 % hors-ligne", "Tout est utilisable sans connexion ; synchro au retour du réseau."),
    ("Installable", "À ajouter sur l'écran d'accueil comme une vraie application."),
    ("Aide intégrée", "Tutoriels consultables à tout moment depuis le menu."),
    ("Évolue avec vous", "Vos retours (corrections, idées) améliorent l'outil en continu."),
]
top = Inches(1.4); bw = Inches(6.0); bh = Inches(1.85)
for i, (n_, d_) in enumerate(items):
    r = i // 2; c = i % 2
    x = Inches(0.6) + (bw + Inches(0.15)) * c
    y = top + (bh + Inches(0.2)) * r
    add_box(s, x, y, bw, bh, n_, d_, title_size=18, body_size=14)
band_bottom(s, "Un dysfonctionnement, une idée ? Remontez-le : l'appli s'améliore en permanence",
            top=Inches(5.55), fill=JAUNE, color=VERT, size=14)
add_footer(s, 13, TOTAL)

# ========================================================================
# Slide 14 — BONNES PRATIQUES ICP
# ========================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Bonnes pratiques — À faire / À éviter")
add_text_box(s, Inches(0.6), Inches(1.2), Inches(6.0), Inches(0.5),
             "À FAIRE", size=20, bold=True, color=VERT)
add_bullets(s, Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.5), [
    "Synchroniser en wifi AVANT de partir sur site.",
    "Réaliser l'ICP avant l'intervention, pas après.",
    "Prendre des photos systématiquement.",
    "Faire signer le site et l'intervenant sur place.",
    "Activer le Mode Express pour les visites simples.",
], size=15, bullet_color=VERT)
add_text_box(s, Inches(6.9), Inches(1.2), Inches(6.0), Inches(0.5),
             "À ÉVITER", size=20, bold=True, color=ROUGE)
add_bullets(s, Inches(6.9), Inches(1.8), Inches(6.1), Inches(4.5), [
    "Intervenir sans inspection préalable des lieux.",
    "Remplir l'ICP de mémoire, au bureau, après coup.",
    "Oublier de faire signer le représentant du site.",
    "Négliger la délimitation de zone et les consignes.",
    "Laisser une ICP en brouillon sans la valider.",
], size=15, bullet_color=ROUGE)
add_footer(s, 14, TOTAL)

# ========================================================================
# Slide 15 — QUESTIONS
# ========================================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.line.fill.background(); bg.fill.solid(); bg.fill.fore_color.rgb = VERT
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.0), SW, Inches(0.08))
b.line.fill.background(); b.fill.solid(); b.fill.fore_color.rgb = JAUNE
add_text_box(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(1.5),
             "Vos questions ?", size=72, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.5), Inches(2.9), Inches(12.3), Inches(0.8),
             "On prend l'appli en main ensemble — et à tout moment après la formation.",
             size=20, color=JAUNE, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(0.6),
             "T. HOCEDEZ", size=26, bold=True, color=BLANC, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(0.6),
             "Responsable QSE — AGRIWATT", size=20, color=BLANC, align=PP_ALIGN.CENTER)
add_text_box(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.5),
             "Merci pour votre attention.", size=18, color=JAUNE, align=PP_ALIGN.CENTER)

# Sauvegarde
out = r"C:\Users\ThibaultHOCEDEZ\Documents\suivipdp\Formation_ICP_Maintenance_2026-06-11.pptx"
prs.save(out)
print(f"OK -> {out}")
print(f"Slides : {len(prs.slides)}")

txt_out = r"C:\Users\ThibaultHOCEDEZ\Documents\suivipdp\Liste_screenshots_ICP_a_prendre.txt"
with open(txt_out, "w", encoding="utf-8") as f:
    f.write("Liste des captures d'écran — Formation Module ICP (Maintenance)\n")
    f.write(f"Présentation du {DATE_PRESENTATION} — T. HOCEDEZ\n")
    f.write("=" * 70 + "\n\n")
    f.write("Application : https://zbpbasv.cluster121.hosting.ovh.net/suivipdp/\n\n")
    f.write("Conseils :\n")
    f.write("- Format conseillé : PNG, largeur 1400 à 1800 px.\n")
    f.write("- Masquer les données client réelles si partage hors AGRIWATT.\n")
    f.write("- Glisser-déposer la capture dans le rectangle gris correspondant.\n\n")
    f.write("-" * 70 + "\n\n")
    for num, desc in SCREENSHOTS:
        f.write(f"Slide {num:>2}  :  {desc}\n\n")
    f.write("-" * 70 + "\n")
    f.write(f"Total : {len(SCREENSHOTS)} captures à prendre.\n")
print(f"Liste captures : {txt_out}")
